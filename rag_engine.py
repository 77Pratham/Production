# rag_engine.py - Privacy-Preserving RAG Engine with encrypted vector storage
import os
import pickle
import faiss
import numpy as np
import hashlib
import json
from typing import Optional, List, Dict, Any
from sentence_transformers import SentenceTransformer
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from cryptography.fernet import Fernet
import logging

logger = logging.getLogger(__name__)

class PrivacyPreservingRAGEngine:
    def __init__(self, kb_path="knowledge_base", index_file="faiss_index.pkl", 
                 encrypted_docs_file="encrypted_docs.pkl", user_context_dir="user_contexts"):
        self.kb_path = kb_path
        self.index_file = index_file
        self.encrypted_docs_file = encrypted_docs_file
        self.user_context_dir = user_context_dir
        self.model = SentenceTransformer("all-MiniLM-L6-v2")
        self.index = None
        self.encrypted_documents = []
        self.document_metadata = []
        self.encryption_key = self._load_or_create_encryption_key()
        self.fernet = Fernet(self.encryption_key)
        
        # Create user context directory
        os.makedirs(self.user_context_dir, exist_ok=True)
    
    def _load_or_create_encryption_key(self) -> bytes:
        """Load or create encryption key for document storage"""
        key_file = "rag_encryption.key"
        if os.path.exists(key_file):
            with open(key_file, "rb") as f:
                return f.read()
        else:
            key = Fernet.generate_key()
            with open(key_file, "wb") as f:
                f.write(key)
            return key
    
    def _encrypt_document(self, text: str) -> bytes:
        """Encrypt document content"""
        return self.fernet.encrypt(text.encode())
    
    def _decrypt_document(self, encrypted_text: bytes) -> str:
        """Decrypt document content"""
        try:
            return self.fernet.decrypt(encrypted_text).decode()
        except Exception as e:
            logger.error(f"Decryption error: {e}")
            return ""
    
    def _read_txt(self, filepath: str) -> str:
        """Read text file"""
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    
    def _read_pdf(self, filepath: str) -> str:
        """Read PDF file"""
        text = ""
        try:
            reader = PdfReader(filepath)
            for page in reader.pages:
                text += page.extract_text() or ""
        except Exception as e:
            logger.error(f"PDF reading error for {filepath}: {e}")
        return text
    
    def _read_docx(self, filepath: str) -> str:
        """Read DOCX file"""
        try:
            doc = docx.Document(filepath)
            return "\n".join([p.text for p in doc.paragraphs])
        except Exception as e:
            logger.error(f"DOCX reading error for {filepath}: {e}")
            return ""
    
    def _read_pptx(self, filepath: str) -> str:
        """Read PPTX file"""
        try:
            prs = Presentation(filepath)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        except Exception as e:
            logger.error(f"PPTX reading error for {filepath}: {e}")
            return ""
    
    def build_index(self):
        """Build encrypted FAISS index from knowledge base"""
        docs = []
        encrypted_docs = []
        metadata = []
        
        if not os.path.exists(self.kb_path):
            logger.warning(f"Knowledge base path {self.kb_path} does not exist")
            return
        
        for file in os.listdir(self.kb_path):
            filepath = os.path.join(self.kb_path, file)
            ext = file.lower().split(".")[-1]
            
            try:
                if ext == "txt":
                    text = self._read_txt(filepath)
                elif ext == "pdf":
                    text = self._read_pdf(filepath)
                elif ext == "docx":
                    text = self._read_docx(filepath)
                elif ext == "pptx":
                    text = self._read_pptx(filepath)
                else:
                    continue
                
                if not text.strip():
                    continue
                
                # Store original text for embedding
                docs.append(text)
                
                # Encrypt and store document
                encrypted_text = self._encrypt_document(text)
                encrypted_docs.append(encrypted_text)
                
                # Store metadata
                doc_metadata = {
                    "filename": file,
                    "size": len(text),
                    "hash": hashlib.sha256(text.encode()).hexdigest()[:16],
                    "type": ext
                }
                metadata.append(doc_metadata)
                
            except Exception as e:
                logger.error(f"Error processing {file}: {e}")
                continue
        
        if not docs:
            logger.warning("No documents found to index")
            return
        
        # Create embeddings
        logger.info(f"Creating embeddings for {len(docs)} documents...")
        embeddings = self.model.encode(docs, convert_to_numpy=True, show_progress_bar=True)
        
        # Build FAISS index
        d = embeddings.shape[1]
        self.index = faiss.IndexFlatL2(d)
        self.index.add(embeddings)
        
        # Store encrypted documents and metadata
        self.encrypted_documents = encrypted_docs
        self.document_metadata = metadata
        
        # Save to disk
        index_data = {
            "index": self.index,
            "encrypted_documents": self.encrypted_documents,
            "metadata": self.document_metadata
        }
        
        with open(self.index_file, "wb") as f:
            pickle.dump(index_data, f)
        
        logger.info(f"Knowledge base indexed with {len(docs)} documents.")
    
    def load_index(self):
        """Load encrypted FAISS index"""
        if not os.path.exists(self.index_file):
            logger.warning("No index found. Run build_index().")
            return
        
        try:
            with open(self.index_file, "rb") as f:
                index_data = pickle.load(f)
            
            self.index = index_data["index"]
            self.encrypted_documents = index_data["encrypted_documents"]
            self.document_metadata = index_data["metadata"]
            
            logger.info("Encrypted FAISS index loaded successfully.")
        except Exception as e:
            logger.error(f"Error loading index: {e}")
    
    def get_secure_context(self, encrypted_query: str, user_id: str, top_k: int = 3) -> Optional[str]:
        """Get context using privacy-preserving retrieval"""
        if not self.index or not self.encrypted_documents:
            return None
        
        try:
            # For demo purposes, decrypt query for embedding
            # In production, use homomorphic encryption
            from security_manager import SecureNLPInterface
            security = SecureNLPInterface()
            query = security.decrypt_data(encrypted_query)
            
            # Create query embedding
            query_embedding = self.model.encode([query], convert_to_numpy=True)
            
            # Search in FAISS index
            distances, indices = self.index.search(query_embedding, top_k)
            
            # Retrieve and decrypt relevant documents
            contexts = []
            for i, idx in enumerate(indices[0]):
                if idx < len(self.encrypted_documents):
                    # Decrypt document
                    decrypted_doc = self._decrypt_document(self.encrypted_documents[idx])
                    if decrypted_doc:
                        # Add user-specific context filtering
                        filtered_context = self._filter_context_for_user(decrypted_doc, user_id)
                        contexts.append({
                            "content": filtered_context,
                            "distance": float(distances[0][i]),
                            "metadata": self.document_metadata[idx]
                        })
            
            # Store user interaction for privacy analysis
            self._log_user_access(user_id, query, [ctx["metadata"]["filename"] for ctx in contexts])
            
            # Combine contexts
            if contexts:
                combined_context = "\n\n".join([ctx["content"] for ctx in contexts[:top_k]])
                return combined_context
            
            return None
            
        except Exception as e:
            logger.error(f"Error in secure context retrieval: {e}")
            return None
    
    def _filter_context_for_user(self, document: str, user_id: str) -> str:
        """Filter document content based on user permissions"""
        # Implement user-specific filtering logic
        # For now, return full document - in production, implement role-based access
        return document
    
    def _log_user_access(self, user_id: str, query: str, accessed_files: List[str]):
        """Log user access for privacy auditing"""
        access_log = {
            "timestamp": datetime.datetime.now().isoformat(),
            "user_id": user_id,
            "query_hash": hashlib.sha256(query.encode()).hexdigest()[:16],
            "accessed_files": accessed_files,
            "query_length": len(query)
        }
        
        user_log_file = os.path.join(self.user_context_dir, f"{user_id}_access.log")
        try:
            with open(user_log_file, "a") as f:
                f.write(json.dumps(access_log) + "\n")
        except Exception as e:
            logger.error(f"Failed to log user access: {e}")
    
    def answer_query(self, query: str, context: Optional[str] = None) -> str:
        """Answer query using retrieved context"""
        if not context:
            return "I couldn't find relevant information in the knowledge base."
        
        # Simple context-based answering
        # In production, use a language model for better responses
        if any(keyword in query.lower() for keyword in ["what is", "define", "explain"]):
            # Extract definition from context
            sentences = context.split('. ')
            for sentence in sentences[:3]:  # First 3 sentences
                if len(sentence) > 20:  # Meaningful sentence
                    return sentence.strip()
        
        # Return first paragraph as answer
        paragraphs = context.split('\n\n')
        if paragraphs:
            return paragraphs[0].strip()[:500] + "..." if len(paragraphs[0]) > 500 else paragraphs[0].strip()
        
        return context[:200] + "..." if len(context) > 200 else context
    
    def get_user_context_summary(self, user_id: str) -> Dict[str, Any]:
        """Get privacy-preserving summary of user interactions"""
        user_log_file = os.path.join(self.user_context_dir, f"{user_id}_access.log")
        
        if not os.path.exists(user_log_file):
            return {"total_queries": 0, "accessed_files": [], "last_access": None}
        
        try:
            queries = 0
            files_accessed = set()
            last_access = None
            
            with open(user_log_file, "r") as f:
                for line in f:
                    log_entry = json.loads(line.strip())
                    queries += 1
                    files_accessed.update(log_entry.get("accessed_files", []))
                    last_access = log_entry.get("timestamp")
            
            return {
                "total_queries": queries,
                "unique_files_accessed": len(files_accessed),
                "accessed_files": list(files_accessed),
                "last_access": last_access
            }
            
        except Exception as e:
            logger.error(f"Error getting user context summary: {e}")
            return {"error": str(e)}
    
    def cleanup_user_data(self, user_id: str, older_than_days: int = 30):
        """Clean up old user data for privacy compliance"""
        user_log_file = os.path.join(self.user_context_dir, f"{user_id}_access.log")
        
        if not os.path.exists(user_log_file):
            return
        
        try:
            cutoff_date = datetime.datetime.now() - datetime.timedelta(days=older_than_days)
            recent_logs = []
            
            with open(user_log_file, "r") as f:
                for line in f:
                    log_entry = json.loads(line.strip())
                    log_date = datetime.datetime.fromisoformat(log_entry["timestamp"])
                    if log_date > cutoff_date:
                        recent_logs.append(line)
            
            # Rewrite file with only recent logs
            with open(user_log_file, "w") as f:
                f.writelines(recent_logs)
            
            logger.info(f"Cleaned up old data for user {user_id}")
            
        except Exception as e:
            logger.error(f"Error cleaning up user data: {e}")

import datetime